import datetime
import math
import os

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt

dfMaster = pd.read_excel('NewFormatMasterPrL2.xlsx', engine='openpyxl')
the_zero = 0
# задачи на открытых колбасках
r = r'(Issue)'
ar = dfMaster.loc[(dfMaster.SLIDE == 'Gantt1')]
obr = ar['VARIABLE'].str.contains(r, na=False)
IsArray1 = ar[obr]

ar2 = dfMaster.loc[(dfMaster.SLIDE == 'Gantt2')]
obr = ar2['VARIABLE'].str.contains(r, na=False)
IsArray2 = ar2[obr]

ar3 = dfMaster.loc[(dfMaster.SLIDE == 'Gantt3')]
obr = ar3['VARIABLE'].str.contains(r, na=False)
IsArray3 = ar3[obr]

ar4 = dfMaster.loc[(dfMaster.SLIDE == 'Gantt4')]
obr = ar4['VARIABLE'].str.contains(r, na=False)
IsArray4 = ar4[obr]

ar5 = dfMaster.loc[(dfMaster.SLIDE == 'Gantt5')]
obr = ar5['VARIABLE'].str.contains(r, na=False)
IsArray5 = ar5[obr]

# таблица с массивами
r = r'(Array)'
ag = dfMaster.loc[(dfMaster.SLIDE == 'Gantt')]
obr = ag['VARIABLE'].str.contains(r, na=False)
# print(len(ag[obr[True]]))
ag = ag[obr]
table_w_arrays = ag
# s = table_w_arrays[5:]
# print(s.iloc[1][3])


# r = r'(Array)'
# ag = dfMaster.loc[(dfMaster.SLIDE == 'Gantt')]
# if


# таблица для каврталов -----------------------------------
r = r'(Quarter)'
ag = dfMaster.loc[(dfMaster.SLIDE == 'All')]
obr = ag['VARIABLE'].str.contains(r, na=False)
ag = ag[obr]
table_w_quartres = ag

# name perc Team таблица для раскрытия ее состовляющих внутри цикла
r = r'(Name_perc_Team)'
ag = dfMaster.loc[(dfMaster.SLIDE == 'Gantt')]
obr = ag['VARIABLE'].str.contains(r, na=False)
ag = ag[obr]
table_w_npt = ag

# table for gantt progress 1st slide

r = r'(Progress_on)'
ag = dfMaster.loc[(dfMaster.SLIDE == 'Gantt')]
obr = ag['VARIABLE'].str.contains(r, na=False)
ag = ag[obr]
g0ponq = ag

header_top = Cm(0.51)
header_left = Cm(0.87)

now = datetime.datetime.now()


# Сложносоставные объекты и фон слайдов
point = 'point1.png'
gannt2 = 'gannt2.png'
oppened = 'oppened.png'
home = 'Home button.png'

pr1 = Presentation()

layouts = pr1.slide_layouts[6]
pr1.slide_width = Cm(33.867)
pr1.slide_height = Cm(19.05)

top = Inches(0)
left = Inches(0)
width = Cm(33.867)
height = Cm(19.05)


def grey_block(shape):
    fillsolid(shape, RGBColor(49, 51, 63))

    tb = shape.text_frame
    adjs = shape.adjustments
    adjs[0] = 0.1

    liner(shape, Pt(3), RGBColor(49, 51, 63))

    shadow = shape.shadow
    shadow.inherit = False


def bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(49, 51, 63)


def value(val):
    df = pd.read_excel('NewFormatMasterPrL2.xlsx', engine='openpyxl')
    rows, cols = np.where(df == val)
    if len(rows) == 0:
        return None
    row, col = rows[0], cols[0]
    if col + 1 > df.shape[1]:
        return None
    return df.iat[row, col + 1]


def value2(val):
    df = pd.read_excel('NewFormatMasterPrL2.xlsx', engine='openpyxl')
    rows, cols = np.where(df == val)
    if len(rows) == 0:
        return None
    row, col = rows[0], cols[0]
    if col + 2 > df.shape[1]:
        return None
    return df.iat[row, col + 2]


def value3(val):
    df = pd.read_excel('NewFormatMasterPrL2.xlsx', engine='openpyxl')
    rows, cols = np.where(df == val)
    if len(rows) == 0:
        return None
    row, col = rows[0], cols[0]
    if col + 3 > df.shape[1]:
        return None
    return df.iat[row, col + 3]


def fillsolid(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def fonter(q, font_size, font_name, font_color=None, alignment=None):
    q.font.size = font_size
    q.font.name = font_name
    if font_color != None:
        q.font.color.rgb = font_color
    if alignment != None:
        q.alignment = alignment


def liner(shape, length, color=None):
    line = shape.line
    line.width = length
    if color != None:
        line.color.rgb = color


def textshape(slide, left, top, width, height, font_color, font_size, font_name, text, paragraphs=None, text0=None,
              align=None, to_slide=None, font_size2=None, font_name2=None, fill_color=None, resize=None, Pt=None,
              lm=None,
              tm=None, bm=None):
    if resize != None:
        if Pt == 1:
            wi = 0.18
        elif Pt == 2:
            wi = 0.42
        elif Pt == 3:
            wi = 0.65
        else:
            raise SyntaxError('Неправильно набран номер для resize')
        width = Cm(len(text) * wi)
    shape = slide.shapes.add_textbox(left, top, width, height)
    txt = shape.text_frame
    q = txt.paragraphs[0]
    q.text = text
    fonter(q, font_size, font_name, font_color)

    if align == 1:
        q.alignment = PP_ALIGN.RIGHT
    elif align == 2:
        q.alignment = PP_ALIGN.CENTER
    if paragraphs == 1:
        if font_size2 == None:
            font_size = font_size
        else:
            font_size = font_size2
        if font_name2 == None:
            font_name = font_name
        else:
            font_name = font_name2
        a = txt.add_paragraph()
        a.text = text0
        fonter(a, font_size, font_name, font_color, PP_ALIGN.LEFT)
        if align == 1:
            a.alignment = PP_ALIGN.RIGHT
        elif align == 2:
            a.alignment = PP_ALIGN.CENTER
    if to_slide != None:
        shape.click_action.target_slide = to_slide
    if fill_color != None:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = fill_color
    if bm != None:
        q.margin_bottom = Cm(bm)
    if tm != None:
        q.margin_top = Cm(top)
    if lm != None:
        q.margin_left = Cm(left)


def long_line(slide, left, top, color, Pt):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Cm(0.004), Cm(17.4))
    fillsolid(shape, color)
    liner(shape, Pt, color)
    shadow = shape.shadow
    shadow.inherit = False

    # slide.shapes._spTree.remove(shape._element)
    # slide.shapes._spTree.insert( slide.shapes._spTree.max_shape_id + 1, shape._element)


def qshape(slide, l0, color=None, name=None, first_one=None):
    qtop = header_top
    width = Cm(2.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l0, qtop, width, Cm(0.6))
    shadow = shape.shadow
    shadow.inherit = False

    if color != None:
        fillsolid(shape, RGBColor(70, 72, 83))
        liner(shape, Pt(0.0003), RGBColor(70, 72, 83))
        color = RGBColor(255, 255, 255)
    else:
        fillsolid(shape, RGBColor(51, 63, 80))
        liner(shape, Pt(0.0003), RGBColor(51, 63, 80))
        color = RGBColor(97, 120, 153)
    if first_one != None:
        fillsolid(shape, RGBColor(70, 72, 83))
        liner(shape, Pt(0.0003), RGBColor(70, 72, 83))
        color = RGBColor(255, 255, 255)

    text = shape.text_frame

    q = text.paragraphs[0]
    q.text = name
    q.font.size = Pt(10.5)
    q.font.name = ('Segoe UI Semibold')
    q.font.color.rgb = color

    long_line(slide, l0 + Cm(2.94), qtop + Cm(0.7), RGBColor(70, 72, 83), Pt(0.00001))

    # shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l0 + Cm(2.7), qtop + Cm(0.7), Cm(0.004), Cm(17.4))
    # fillsolid(shape, RGBColor(37,64,97))
    # liner(shape, Pt(0.00001), RGBColor(37,64,97))
    # shadow = shape.shadow
    # shadow.inherit = False
    #
    # slide.shapes._spTree.remove(shape._element)
    # slide.shapes._spTree.insert(2, shape._element)
    if slide == slideGannt:
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

def Quarters(slide, array, gonq):
    if len(array) % 3 != 0:
        raise ValueError('Wrong array length. Make it dividable by 3')
    mod3 = len(array) // 3
    lean = 22.51
    margin = Cm(lean / mod3)
    first = ((array[3:].index(1) // 3) * margin) + margin + Cm(0.6)
    mc = math.ceil(array.count(1) / 3)
    count = mc * (margin)
    qtop = header_top + Cm(0.8)
    long_line(slide, Cm(4.0), header_top + Cm(0.7), RGBColor(70, 72, 83), Pt(0.00001))

    if slide == slideGannt or slide == slideGannt2:
        array = [1] * len(array)
        mod3 = len(array) // 3
        margin = Cm(lean / mod3)

    else:
        mod3 = len(array) // 3
        lean = 22.51
        margin = Cm(lean / mod3)
        first = ((array.index(1) // 3) * margin) + margin + Cm(0.6)
        textshape(slide, first - Cm(1.8), qtop - Cm(0.21), Cm(2), Cm(1.09), RGBColor(255, 255, 255), Pt(9),
                  'Segoe UI Black', 'Progress')
    l0 = Cm(4.36)
    # if slide != slideGannt:
    for i in range(mod3):
        # print(table_w_quartres)
        name = table_w_quartres.iloc[i][3]
        if 1 in array[:3]:
            qshape(slide, l0, 1, name)
            if slide != slideGannt and slide != slideGannt2:
                onq = int(gonq.iloc[i][3])
                height_onq = Cm(onq / 100 * 2.6)
                rectangler(slide, l0, qtop, height_onq, Cm(0.3), 2, Pt(0.3), RGBColor(49, 51, 68))
                textshape(slide, l0 + Cm(2.0), qtop - Cm(0.15), Cm(0.7), Cm(0.6), RGBColor(255, 255, 255), Pt(9),
                          'Segoe UI Black', f'{onq}')
                textshape(slide, l0 + Cm(2.5), qtop - Cm(0.1), Cm(0.4), Cm(0.5), RGBColor(255, 255, 255), Pt(6),
                          'Segoe UI Black', '%')
            l0 = l0 + margin

            # first = l0
        else:
            qshape(slide, l0, None, name)
            l0 = l0 + margin
        array = array[3:]


def now_line(month, slide):
    for i in range(12):
        lni = month * 1.6625
        left = (Cm(5.39) + (Cm(lni)))
        if i == month:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Cm(1.6), Cm(0.06), Cm(16.4))
            fillsolid(shape, RGBColor(192, 0, 0))
            liner(shape, Pt(0.003), RGBColor(192, 0, 0))
            shadow = shape.shadow
            shadow.inherit = False
            textshape(slide, left + Cm(0.4), Cm(1.6), Cm(0.58), Cm(0.71), RGBColor(192, 0, 0), Pt(7), 'Segoe UI Black',
                      'NOW')


def oval(slide, color, left, top):
    width = Cm(0.21)
    height = Cm(0.21)

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)

    fillsolid(shape, color)
    liner(shape, Pt(0.01), RGBColor(49, 51, 63))

    shadow = shape.shadow
    shadow.inherit = False


# def home_button(slide):
#     shape = slide.shapes.add_picture(home, Cm(31.65), Cm(0.82), Cm(1.40), Cm(1.40))
#     shape.click_action.target_slide = slide1
#     shadow = shape.shadow
#     shadow.inherit = False


def SberGradientLineAdjus(shape):
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 180

    gr_stop = fill.gradient_stops
    gr_stop1 = gr_stop[-1]
    gr_stop2 = gr_stop[0]
    gr_stop3 = gr_stop[1]
    color1 = gr_stop1.color
    color1.rgb = RGBColor(84, 130, 53)
    color2 = gr_stop2.color
    color2.rgb = RGBColor(84, 130, 53)
    color3 = gr_stop3.color
    color3.rgb = RGBColor(94, 158, 239)
    adjs = shape.adjustments
    adjs[0] = 0.2

    liner(shape, Pt(6), RGBColor(49, 51, 63))

    shadow = shape.shadow
    shadow.inherit = False


def bigger_flowchart_shape(color, left, top, length, slide, name2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=left, top=top, width=length, height=Cm(0.43))
    # shape3.rotation = 90.0

    fillsolid(shape, color)

    liner(shape, Pt(0.75), color)

    shadow = shape.shadow
    shadow.inherit = False

    textshape(slide, length + Cm(0.2) + left, top - Cm(0.1), Cm(5), Cm(0.43), RGBColor(255, 255, 255), Pt(8),
              'Segoe UI Light', f'{name2}', None, None, None, None, None, None, RGBColor(49, 51, 63), 1, 1)
    # RGBColor((49, 51, 63))


def not_selected(slide, to_slide, left, top, perc, txt2, txt3):
    perc = f'{perc}'
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=left - Cm(4.72), top=top + Cm(0.37),
                                   width=Cm(0.81), height=Cm(0.47))

    fillsolid(shape, RGBColor(70, 72, 83))

    liner(shape, Pt(1), RGBColor(94, 158, 239))

    # Плашка с процентами
    textshape(slide, left - Cm(4.78), top + Cm(0.33), Cm(1.01), Cm(0.56), RGBColor(255, 255, 255), Pt(7),
              'Segoe UI Semibold', perc, None, None, None, to_slide)
    textshape(slide, left - Cm(4.5), top + Cm(0.4), Cm(1.01), Cm(0.56), RGBColor(255, 255, 255), Pt(5),
              'Segoe UI Semibold', '%', None, None, None, to_slide)
    textshape(slide, left - Cm(5.02), top + Cm(0.93), Cm(4.07), Cm(1.11), RGBColor(255, 255, 255), Pt(10.5),
              'Segoe UI Semibold', txt2, 1, txt3, None, to_slide, Pt(9), 'Segoe UI Light')

    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=left + Cm(23.72), top=top - Cm(0.13),
                                    width=Cm(2.11), height=Cm(2.11))
    liner(shape2, Pt(3), RGBColor(94, 158, 239))

    fillsolid(shape2, RGBColor(70, 72, 83))

    shape2.click_action.target_slide = to_slide

    image = slide.shapes.add_picture(point, left=left + Cm(24.45), top=top + Cm(0.61), width=Cm(0.64), height=Cm(0.56))
    image.click_action.target_slide = to_slide


def not_selected2(slide, to_slide, left, top, perc, txt2, txt3):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=left - Cm(5.82), top=top + Cm(0.07),
                                   width=Cm(3.78), height=Cm(1.91))

    fillsolid(shape, RGBColor(94, 158, 239))

    liner(shape, Pt(1), RGBColor(0, 0, 0))

    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=left - Cm(4.72), top=top + Cm(0.37),
                                    width=Cm(0.81), height=Cm(0.47))
    fillsolid(shape1, RGBColor(255, 255, 255))

    liner(shape1, Pt(1), RGBColor(0, 0, 0))

    # проценты, фамилия и название кластера
    textshape(slide, left - Cm(4.78), top + Cm(0.33), Cm(1.01), Cm(0.56), RGBColor(0, 0, 0), Pt(7), 'Segoe UI Semibold',
              f'{perc}', None, None, None, to_slide)
    textshape(slide, left - Cm(4.5), top + Cm(0.4), Cm(1.01), Cm(0.56), RGBColor(0, 0, 0), Pt(5), 'Segoe UI Semibold',
              '%', None, None, None, to_slide)
    textshape(slide, left - Cm(5.02), top + Cm(0.93), Cm(4.07), Cm(1.11), RGBColor(0, 0, 0), Pt(10.5),
              'Segoe UI Semibold', f'{txt2}', 1, f'{txt3}', None, to_slide, Pt(9), 'Segoe UI Light')

    image = slide.shapes.add_picture(oppened, left=left + Cm(23.72), top=top - Cm(0.13), width=Cm(2.11),
                                     height=Cm(2.11))
    image.click_action.target_slide = to_slide


def greyFill(shape):
    fillsolid(shape, RGBColor(70, 72, 83))

    tb = shape.text_frame
    adjs = shape.adjustments
    adjs[0] = 0.13

    liner(shape, Pt(6), RGBColor(49, 51, 63))

    shadow = shape.shadow
    shadow.inherit = False


def selected_block(shape):
    fillsolid(shape, RGBColor(49, 51, 63))

    tb = shape.text_frame
    adjs = shape.adjustments
    adjs[0] = 0.1

    liner(shape, Pt(3), RGBColor(94, 158, 239))

    shadow = shape.shadow
    shadow.inherit = False


def rectangler(slide, left, top, width, height, color=None, line=None, line_color=None, adj=None, text=None,
               font_color=None, font_size=None, font_name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if color == 2:
        fill = shape.fill
        fill.patterned()
        fill.pattern = MSO_PATTERN.DARK_VERTICAL
        fill.fore_color.rgb = RGBColor(94, 158, 239)
        fill.back_color.rgb = RGBColor(49, 51, 68)

    else:
        fillsolid(shape, color)

    if line != None:
        liner(shape, line, line_color)
        # line = shape.line
        # line.width = line
        # line.fore_color.rgb = line_color

    if adj != None:
        adjs = shape.adjustments
        adjs[0] = adj
    if text != None:
        txt = shape.text_frame
        q = txt.paragraphs[0]
        q.text = f'{text}'
        fonter(q, font_size, font_name, font_color)

    shadow = shape.shadow
    shadow.inherit = False


def flowchart_shape(color, top, slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR, left=Cm(4.31), top=top, width=Cm(0.34),
                                   height=Cm(0.23))
    shape.rotation = -90.0

    fillsolid(shape, color)

    liner(shape, Pt(0.75), color)

    shadow = shape.shadow
    shadow.inherit = False
    shape.rotation = -90.0



def hex_to_rgb(value):
    # value = value.lstrip('#')
    lv = len(value)
    return tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3))



def Sausage2(slide, slideNumber, top, array, to_slide, name, surname, perc, l, array2, lengthA, gponq, i):

    ni = array.index(1)
    # print(ni)
    # Quarters(slide, array, gponq)
    nj = array.count(1)
    length = 19.95 / len(array)
    lni = length * ni
    if i>=5:
        l = l+5
        to_slide1 = slideGannt2
    else:
        to_slide1 = slideGannt


    if l != slideNumber - 1:
        left = Cm(5.39 + lni)
        width = Cm(length * nj)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height=Cm(2.26))
        greyFill(shape)
        shape.click_action.target_slide = to_slide

        not_selected(slide, to_slide, Cm(5.39), top, perc, name, surname)
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(3, shape._element)

    else:
        Quarters(slide, array, gponq)

        not_selected2(slide, to_slide1, Cm(5.39), top, perc, name, surname)
        # проверяем количество подзадач (если больше двух, то увеличиваем их высоту)
        if lengthA > 2:
            # номера слайдов - 4 или 5,то поднимаем блоки вверх
            if (l == 4 or l == 3 or l == 8 or l == 9):
                remain = lengthA - 2
                height = Cm(0.7) * remain
                top = top - height + Cm(1.4)

            height = lengthA * Cm(0.7)

            width = Cm(length * nj)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.25 + lni), top=top - Cm(0.14),
                                           width=width + Cm(0.29), height=height + Cm(0.26))
            grey_block(shape)

            # slide.shapes._spTree.remove(shape._element)
            # slide.shapes._spTree.insert(2, shape._element)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.39 + lni), top=top,
                                           width=width, height=height)
            selected_block(shape)

            # slide.shapes._spTree.remove(shape._element)
            # slide.shapes._spTree.insert(3, shape._element)

        else:
            width = Cm(length * nj)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.25 + lni), top=top - Cm(0.14),
                                           width=width + Cm(0.29), height=Cm(2.52))
            grey_block(shape)
            slide.shapes._spTree.remove(shape._element)
            slide.shapes._spTree.insert(2, shape._element)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.39 + lni), top=top,
                                           width=width, height=Cm(2.26))
            selected_block(shape)

            # slide.shapes._spTree.remove(shape._element)
            # slide.shapes._spTree.insert(3, shape._element)
        #     width = Cm(length * nj)
        #
        #     shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.25 + lni), top=top - Cm(0.14),
        #                                    width=width + Cm(0.29), height=lengthA*Cm(0.7) + Cm(0.26))
        #     grey_block(shape)
        #     shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Cm(5.39 + lni), top=top,
        #                                    width=width, height=lengthA * Cm(0.7))
        #     selected_block(shape)
        for h in range(lengthA):
            # dfcs = {1: RGBColor(251, 243, 42), 2: (0, 176, 240), 3: slide5Gannt, 4: (160, 208, 123), 5: (175, 171, 171), 6:(244, 102, 119)}
            if lengthA>16:
                raise ValueError('слишком много подзадач. Проверьте Excel файл')

            # if clr not in dfcs.keys():
            #     raise ValueError('Неправильный номер')
            # else:
            #     print(dfcs[clr])
            #     color = f'RGBColor{dfcs[clr]}'

            clr = array2.iloc[h][5]

            if clr == 1:
                color = RGBColor(251, 243, 42)
            elif clr == 2:
                color = RGBColor(0, 176, 240)
            elif clr == 3:
                color = RGBColor(47, 85, 151)
            elif clr == 4:
                color = RGBColor(160, 208, 123)
            elif clr == 5:
                color = RGBColor(175, 171, 171)
            elif clr == 6:
                color = RGBColor(244, 102, 119)

            if h == 0:
                he = lengthA*Cm(0.63)
                # ptop = Cm(top) + height -lengthA * Cm(0.61)
                top = top + (height -(lengthA * Cm(0.63))+Cm(0.1))/2





            beg = int(array2.iloc[h][3])
            # if beg>ni or beg>(ni+nj):
            #     raise ValueError('The issue is out of range!')

            leng2 = length / 2
            name2 = array2.iloc[h][6]
            length_1 = Cm((int(array2.iloc[h][4]) - int(array2.iloc[h][3]) + 1) * (length))+Cm(0.01)  # длина подзадачи
            start_point = Cm(beg * leng2 + 5.60 + lni)  # начало подзадачи
            # if beg<ni:
            #     raise ValueError('Начало меньше конца')

            if length_1>width:
                raise ValueError('слишком большая подзадача')


            bigger_flowchart_shape(color, start_point, top, length_1, slide, name2)
            top = top + Cm(0.65)
            # ptop = ptop + Cm(0.65)


def Sausage(array, to_slide, top, tuple, name, surname, perc, sl = None):
    if sl == None:
        sl = slideGannt
    # elif sl == slideGannt2:
    #     top =
    for i in range(len(array)):
        ni = array.index(1)
        nj = array.count(1)
        # 19.95 - это максимальная длина колбасы
        length = 19.95 / len(array)
        lni = length * ni
        # 1.6625 - столько весит один месяц при общей отрисовке на 12
        # if len(array) / 3 == 7:
        #     lni = ni * 1.6625 * 4 / 7
        #     length = 1.6625 * 4 / 7
        # elif len(array) / 3 == 4:
        #     lni = ni * 1.6625 * 2 / 3
        #     length = 1.6625 * 2 / 3
        if i == ni:
            # добавляем левый отсутп к началу массива
            left = (Cm(5.39) + (Cm(lni)))
            width = Cm(length * nj)

            shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height=Cm(2.26))
            SberGradientLineAdjus(shape)
            shape.click_action.target_slide = to_slide

            not_selected(sl, to_slide, Cm(5.39), top, perc, name, surname)
            for k in range(4):
                top1 = top + Cm(0.4)
                color = int(tuple.iloc[k][5])
                ni = int(tuple.iloc[k][3])
                text = tuple.iloc[k][4]
                lni = ni * length
                left = (Cm(5.69) + (Cm(lni)))

                # if nik>ni:
                #     raise ValueError('The issue is out of range!')

                shape2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, left, top=top + Cm(0.2), width=Cm(0.30),
                                             height=Cm(0.30))
                if color == 1:
                    fillsolid(shape2, RGBColor(251, 243, 42))
                    liner(shape2, Pt(0.003), RGBColor(251, 243, 42))

                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)
                if color == 2:
                    fillsolid(shape2, RGBColor(0, 176, 240))
                    liner(shape2, Pt(0.003), RGBColor(0, 176, 240))

                    flowchart_shape(RGBColor(0, 176, 240), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)
                if color == 3:
                    fillsolid(shape2, RGBColor(47, 85, 151))

                    liner(shape2, Pt(0.003), RGBColor(47, 85, 151))
                    flowchart_shape(RGBColor(47, 85, 151), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)
                if color == 4:
                    fillsolid(shape2, RGBColor(160, 208, 123))

                    liner(shape2, Pt(0.003), RGBColor(160, 208, 123))
                    flowchart_shape(RGBColor(160, 208, 123), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)
                if color == 5:
                    fillsolid(shape2, RGBColor(175, 171, 171))
                    liner(shape2, Pt(0.003), RGBColor(175, 171, 171))

                    flowchart_shape(RGBColor(175, 171, 171), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)
                if color == 6:
                    fillsolid(shape2, RGBColor(244, 102, 119))
                    liner(shape2, Pt(0.003), RGBColor(244, 102, 119))

                    flowchart_shape(RGBColor(244, 102, 119), top1, slideGannt)
                    flowchart_shape(RGBColor(251, 243, 42), top1, slideGannt2)

                    top1 = top1 + Cm(0.45)

                shadow2 = shape2.shadow
                shadow2.inherit = False

                top = top + Cm(0.46)
                shape3 = sl.shapes.add_textbox(left=left + Cm(0.13), top=top - Cm(0.38), width=Cm(11.12),
                                               height=Cm(0.6))
                txt = shape3.text_frame
                q = txt.paragraphs[0]
                q.text = text
                fonter(q, Pt(7), 'Segoe UI Light', RGBColor(255, 255, 255))


def text_frames(text, slide, left, width):
    rectangler(slide, left, Cm(3.06), width, Cm(0.68), RGBColor(48, 51, 63), Pt(0.0001), RGBColor(48, 51, 63), None,
               text, RGBColor(255, 255, 255), Pt(10), 'Segoe UI Semibold')
    # shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Cm(3.06), width, Cm(0.68))
    # fillsolid(shape, RGBColor(48, 51, 63))
    # liner(shape, Pt(0.0001), RGBColor(48, 51, 63))
    #
    #
    #
    # txt = shape.text_frame
    # q = txt.paragraphs[0]
    # q.text = f'{text}'
    # fonter(q, Pt(10), 'Segoe UI Semibold', RGBColor(255, 255, 255))


def circle(slide, left, top, radius, color, Pt, line_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, radius, radius)
    fillsolid(shape, color)
    liner(shape, Pt, line_color)

    shadow = shape.shadow
    shadow.inherit = False


layouts = pr1.slide_layouts[6]
slideGannt = pr1.slides.add_slide(layouts)
bg(slideGannt)
# shapeGannt = slideGannt.shapes.add_picture(gannt2, left, top, width, height)
# home_button(slideGannt)
slideGannt2 = pr1.slides.add_slide(layouts)
bg(slideGannt2)
# shapeGannt2 = slideGannt2.shapes.add_picture(gannt2, left, top, width, height)

slide3Gannt = pr1.slides.add_slide(layouts)
bg(slide3Gannt)

slide4Gannt = pr1.slides.add_slide(layouts)
bg(slide4Gannt)

slide5Gannt = pr1.slides.add_slide(layouts)
bg(slide5Gannt)

slide6Gannt = pr1.slides.add_slide(layouts)
bg(slide6Gannt)

slide7Gannt = pr1.slides.add_slide(layouts)
bg(slide7Gannt)

slide8Gannt = pr1.slides.add_slide(layouts)
bg(slide8Gannt)

slide9Gannt = pr1.slides.add_slide(layouts)
bg(slide9Gannt)

slide10Gannt = pr1.slides.add_slide(layouts)
bg(slide10Gannt)

slide11Gannt = pr1.slides.add_slide(layouts)
bg(slide11Gannt)

slide12Gannt = pr1.slides.add_slide(layouts)
bg(slide12Gannt)

slide13Gannt = pr1.slides.add_slide(layouts)
bg(slide13Gannt)

slide14Gannt = pr1.slides.add_slide(layouts)
bg(slide14Gannt)

slide15Gannt = pr1.slides.add_slide(layouts)
bg(slide15Gannt)


now_month = now.month
now_line(now_month, slideGannt)


n = len(table_w_arrays)
df2s = {0: slide3Gannt, 1: slide4Gannt, 2: slide5Gannt, 3: slide6Gannt, 4: slide7Gannt,
        5: slide8Gannt, 6: slide9Gannt, 7: slide10Gannt,
8: slide11Gannt, 9: slide12Gannt,10: slide13Gannt, 11: slide14Gannt, 12: slide15Gannt}

aloc = dfMaster.loc[(dfMaster.SLIDE == 'Gantt') & (dfMaster.VARIABLE == f'Array_{1}')]
aloc = aloc.iloc[0][3]
aloc = list(str(aloc))
array = [int(item) for item in aloc]

Quarters(slideGannt,array, None)
Quarters(slideGannt2, array, None)


for i in range(n):
    # print(table_w_arrays)
    if i>=5:
        sl = slideGannt2
        top = Cm(3.08) + (i-5) * (Cm(3.0))
        # m =
    else:
        sl = slideGannt
        # топ для колбасы
        top = Cm(3.08) + (i) * (Cm(3.0))


    #фильтр по задачам Gantt
    r = r'(Issue)'
    ar = dfMaster.loc[(dfMaster.SLIDE == f'Gantt{i + 1}')]
    obr = ar['VARIABLE'].str.contains(r, na=False)
    IsArray = ar[obr]
    array2 = IsArray


    #фильтр по процентам прогресса Gantt
    r = r'(Progress_on)'
    ag = dfMaster.loc[(dfMaster.SLIDE == f'Gantt{i + 1}')]
    obr = ag['VARIABLE'].str.contains(r, na=False)
    ag = ag[obr]
    gponq = ag

    r = f'Name_perc_Team_{i + 1}'
    ag = dfMaster.loc[(dfMaster.SLIDE == f'Gantt')]
    # print(ag)
    obr = ag['VARIABLE'].str.contains(r, na=False)
    ag = ag[obr]
    npc = ag


    r = f'Issue{i + 1}'
    ag = dfMaster.loc[(dfMaster.SLIDE == f'Gantt')]
    obr = ag['VARIABLE'].str.contains(r, na=False)
    ag = ag[obr]
    tuple = ag

    if i==5:
        print(tuple)


    surname = npc.iloc[0][3]
    perc = npc.iloc[0][4]
    name = npc.iloc[0][5]

    #фильтр и проебразование массива для первичного слайда Gantt
    aloc = dfMaster.loc[(dfMaster.SLIDE == 'Gantt') & (dfMaster.VARIABLE == f'Array_{i+1}')]

    aloc = aloc.iloc[0][3]
    aloc = list(str(aloc))
    array = [int(item) for item in aloc]

    if i not in df2s.keys():
        # exec('slide%Gantt = pr1.slides.add_slide(layouts) bg(slide%Gantt)'%i)
        # bg(slide8Gannt)
        raise ValueError('Добавить слайды в код')
    else:
        to_slide1 = df2s[i]

    Sausage(array, to_slide1, top, tuple, name, surname, perc, sl)

    lengthA = value(f'Amount_of_issuesGantt{i+1}')

    # Условие на количество Предметных областей
    # Актуально при количестве предметных областей 5<...<10
    # Будет отрисовываться вычисленное число колбасок
    k = None
    if i>=5:
        le = n-5
        sl = slideGannt2


    else:
        sl = slideGannt
        le = 5


    for l in range(le):
        k = None

        if i>=5:
            k = True

            # чтобы фильтр с именами подстроилась под
            r = f'Name_perc_Team_{l + 6}'
            ag = dfMaster.loc[(dfMaster.SLIDE == f'Gantt')]
            obr = ag['VARIABLE'].str.contains(r, na=False)
            ag = ag[obr]
            npc = ag




        else:
            r = f'Name_perc_Team_{l + 1}'
            ag = dfMaster.loc[(dfMaster.SLIDE == f'Gantt')]
            obr = ag['VARIABLE'].str.contains(r, na=False)
            ag = ag[obr]
            npc = ag

        # if i == 2 and l == 1:
        #     print(npc.iloc[0][3])

        # топ для колбасы
        top = Cm(3.08) + (l) * (Cm(3.0))


        # print(table_w_arrays.iloc[kl][3])
        aloc = table_w_arrays.iloc[l][3].replace(',', '')
        aloc = list(str(aloc))
        array = [int(item) for item in aloc]

        surname = npc.iloc[0][3]
        perc = npc.iloc[0][4]
        name = npc.iloc[0][5]

        if i not in df2s.keys():
            raise ValueError('Добавить слайды в код для второго цикла')
        else:
            to_slide2 = df2s[l]

        # Условие для замены переменных для заполнения массива
        if k == True:
            # table_w_arrays = s
            surname = npc.iloc[0][3]
            perc = npc.iloc[0][4]
            name = npc.iloc[0][5]
            to_slide2 = df2s[l+5]
            aloc = table_w_arrays.iloc[l+5][3].replace(',', '')
            aloc = list(str(aloc))
            array = [int(item) for item in aloc]


        Sausage2(to_slide1, i + 1, top, array, to_slide2, name, surname, perc, l, array2, lengthA, gponq, i)
        the_zero +=1



pr1.save('Gantt.pptx')
os.startfile('Gantt.pptx')
